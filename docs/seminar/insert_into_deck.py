"""Augment the seminar deck with the tensor-view figure + new slides.

Reads ``HyMeKo_Seminar.pptx`` and writes a *new* ``HyMeKo_Seminar.with_refs.pptx``
(the original is never overwritten — it is the user's artifact). It:

  1. places the clean *aggregated* tensor-view figure (figures/tensor_view.png)
     on the "Star vs clique" slide — imported, not a faded background;
  2. adds "References (selected)", "Publications & submissions" and "Future work"
     slides;
  3. reorders so the "Conclusion & outlook" slide is LAST.

Run:  python docs/seminar/insert_into_deck.py
Requires python-pptx (user-approved dev tooling, 2026-06-15: APPROVED-CORE-EDIT:
pptx-seminar-tooling).
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

HERE = Path(__file__).resolve().parent
FIG = HERE / "figures"
SRC = HERE / "HyMeKo_Seminar.pptx"
DST = HERE / "HyMeKo_Seminar.with_refs.pptx"

HEADER_FONT = "Century Gothic"   # Futura's Windows-native fallback (deck note)
BODY_FONT = "Bahnschrift"
ACCENT = RGBColor(0x4B, 0x2E, 0x83)

REFERENCES = [
    "Berge, Graphs and Hypergraphs, 1973 · Levi, Finite geometrical systems, 1942",
    "Feng et al., Hypergraph Neural Networks (HGNN), AAAI 2019 · Bai et al., Hypergraph Conv/Attention, 2021",
    "Cartwright & Harary, Structural balance, 1956 (σ-cycle balance)",
    "Derr et al., Signed GCN (SGCN), ICDM 2018 · Huang et al., SiGAT, 2019 · SDGNN",
    "Kumar et al., Bitcoin signed nets · Leskovec et al., Epinions/Slashdot",
    "Kolmogorov, 1957 · Liu et al., KAN: Kolmogorov–Arnold Networks, 2024",
    "Forman, combinatorial Ricci curvature, 2003 · Topping et al., SDRF, ICLR 2022 · Hodge decomposition",
    "Quigley et al., ROS/URDF, 2009 · Macenski et al., ROS 2, 2022 · Todorov et al., MuJoCo, 2012 · Gazebo, 2004",
    "OMG SysML v2 · Friedenthal et al., A Practical Guide to SysML · Fey & Lenssen, PyTorch Geometric, 2019",
    "O'Connor et al., BLAKE3, 2020 · Eclipse iceoryx2 · VIATRA · EMF-IncQuery · ATL",
    "Hajdu et al., tensor/generative hypergraph papers (2022), CogInfoCom (2024), HSMM (2026)",
    "SignedKAN / HSiKAN / mixed-arity αₖ line · Friedler et al., P-graph axioms",
]

# Publication portfolio (hand-curated from the paper/ tree, 2026-06-15; confirm
# statuses before presenting). (line, status-tag for the leading marker colour).
PUBLICATIONS = [
    ("HyMeKo: Canonical-Hypergraph Infrastructure for CPS engineering — IEEE SMC 2026 (conf)", "sub"),
    ("HSiKAN: Mixed-Arity Signed Kolmogorov–Arnold Networks — IEEE SMC 2026 (WIP)", "sub"),
    ("AC-HSiKAN: signed-link prediction — Elsevier journal (under review)", "sub"),
    ("Cycles vs. Walks: two inductive biases in Signed KAN — IEEE SISY 2026", "sub"),
    ("Real-time hypergraph contextual layer for robot control (UR5e) — IEEE SISY 2026", "sub"),
    ("Vulkan/WebGPU compute kernels for hypergraph visualisation — GrafGeo 2026", "sub"),
    ("HyMeKo (journal extension) — IEEE T-SMC: Systems", "prep"),
    ("Spectral-entropy regularization, Lyapunov schedule — IEEE T-SMC: Cybernetics", "prep"),
    ("Signed-hypergraph structural priors for relational learning — IEEE TPAMI", "prep"),
    ("HSiKAN + honest σ-masked protocol for signed-graph link prediction — Nature Communications", "prep"),
    ("GPU-kernels journal extension — IEEE T-SMC: Systems", "prep"),
    ("HyMeKo live demo — MDPI Technologies", "plan"),
]
_PUB_PREFIX = {"sub": "[submitted] ", "prep": "[in prep] ", "plan": "[planned] "}

FUTURE_WORK = [
    "σ-masked strict protocol — finish the sign-aware leakage audit + 5-seed grid; lock the honest operating point.",
    "Round-trip .hymeko → torch.nn — structural parity holds; next: runnable round-trip + faithful Soma vision (Hodge/stim/patch).",
    "Broader transform targets — task emitters (BehaviorTree.CPP / PDDL / ROS 2 actions) for the .hymeko task layer.",
    "Larger-scale corpora — bigger signed-graph & vision datasets under the 16 GB RSS discipline.",
    "HSMM → FPGA — Nagare dataflow → HSMM abstract machine → Zynq (theory → systems → compiler).",
    "Authoring surface — in-browser editor: multi-file imports & profiles, parametric generators, arc-value editing.",
]

# --- Kato seminar-review reframing (kato_seminar_review, 2026-06-16) ----------
# The deck makes many strong claims at once; Kato will ask "what is THE central
# contribution / what is proven vs prototype vs program". So state the two coupled
# contributions up front, give a four-block arc, and prep demo / collaboration /
# Q&A slides. Framing one-liners that tame the "dragon heads" (Clifford / HyMeYOLO)
# and gloss "Gömb" live in the presenter notes on the contribution slide.

# One consolidated front slide: the two coupled contributions + the four-block
# arc + the "Gömb" gloss (the gloss the presenter must say first). Merged to keep
# the deck tight per the user's "compress the slide count" note.
CONTRIB_ARC = [
    "TWO coupled contributions, from ONE representation:",
    "(1) Canonical hypergraph infrastructure — a DSL + IR; one source, many targets.",
    "(2) Structural-prior learning — cycles/walks as inductive features (SignedKAN / HSiKAN / Gömb), strictly leakage-audited.",
    "The SAME IR is the engineering transform AND the learning substrate. Frame: a unified substrate, not a bag of unrelated tricks.",
    "Arc:  1) Problem (n-ary; pairwise loses identity) → 2) Framework (DSL+IR+star+codegen) → 3) Learning (cycle/walk prior) → 4) Evidence (honest protocol, results, demo).",
    "Note: 'Gömb' = sphere (Hungarian) — the name of the strict cascade configuration.",
]
# One consolidated closing slide: the minimum demo + the collaboration ask.
DEMO_COLLAB = [
    "Minimum working demo: .hymeko → IR → star expansion → PyTorch / URDF; plus a HSiKAN graph-only kinematics demo.",
    "One source, many targets: URDF / SDF / MJCF / SysML / PyTorch from a single IR.",
    "Efficiency anchor: star vs clique = 1,498 vs 10,991 NNZ; 79.7 ms vs 496 ms.",
    "Collaboration: robot-structure learning — encode morphology + task in HyMeKo; test whether structural priors improve sample efficiency / generalization in control.",
]
QA = [
    "Minimum working demo?  →  .hymeko → IR → star expansion → PyTorch / URDF; plus HSiKAN graph-only kinematics.",
    "Novel vs existing hypergraph NNs?  →  not just an HGNN layer — canonical hypergraph INFRASTRUCTURE + signed cycle/walk tuple prior + strict leakage-audited protocol.",
    "Relation: engineering framework vs learning model?  →  the same IR provides both the engineering transform and the learning substrate.",
    "Is the 0.996 result leakage-free?  →  that is the transductive-convention number, NOT the headline; the strict baseline is the honest claim.",
    "What can we collaborate on?  →  robot-structure learning (morphology + task in HyMeKo; structural priors for control).",
]
def slide_text(slide) -> str:  # type: ignore[no-untyped-def]
    out = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            out.append(sh.text_frame.text)
    return " ".join(out)


def _send_to_back(slide, shape) -> None:  # type: ignore[no-untyped-def]
    """Move a shape to the back of the z-order (drawn first → behind text)."""
    tree = slide.shapes._spTree
    sp = shape._element
    tree.remove(sp)
    # Children 0,1 are <p:nvGrpSpPr>,<p:grpSpPr>; insert right after them.
    tree.insert(2, sp)


def add_tensor_figure(prs) -> bool:  # type: ignore[no-untyped-def]
    """Place the aggregated tensor-view figure (square matrices) as a faded,
    full-width BACKGROUND on the Star-vs-clique slide — sent to the back so the
    slide text reads on top. Uses the more-visible ``tensor_view_bg.png``."""
    fig = FIG / "tensor_view_bg.png"
    if not fig.exists():
        return False
    W, H = prs.slide_width, prs.slide_height
    for slide in prs.slides:
        if "star vs clique" in slide_text(slide).lower():
            # Full-width band, vertically centred; preserve aspect (height auto).
            pic = slide.shapes.add_picture(str(fig), Emu(0), Emu(0), width=W)
            pic.top = Emu(int((H - pic.height) / 2))
            _send_to_back(slide, pic)
            return True
    return False


def _render_cycle_spectrum(chart, out_path: str) -> bool:  # type: ignore[no-untyped-def]
    """Render a clustered-column chart's data as a compact annotated heatmap —
    a "cycle-spectrum fingerprint": rows = mechanisms, columns = cycle length k,
    each cell = number of k-cycles. Far less sparse-looking than grouped bars when
    the data is mostly zeros, and on-theme (a structural-signature matrix). Reads
    the data straight off the embedded chart. Returns False for unsupported types.
    """
    from pptx.enum.chart import XL_CHART_TYPE

    if chart.chart_type != XL_CHART_TYPE.COLUMN_CLUSTERED:
        return False
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    plot = chart.plots[0]
    cats = [str(c) for c in plot.categories]
    series = [(s.name, [float(v or 0.0) for v in s.values]) for s in plot.series]
    rows = [name for name, _ in series]
    mat = np.array([vals for _, vals in series], dtype=float)  # (rows, cols)
    if mat.size == 0:
        return False
    vmax = max(float(mat.max()), 1.0)

    fig, ax = plt.subplots(figsize=(8.2, 2.58), dpi=200)
    ax.imshow(mat, cmap="Purples", aspect="auto", vmin=0.0, vmax=vmax)
    ax.set_xticks(range(len(cats)))
    ax.set_xticklabels(cats, fontsize=11)
    ax.set_yticks(range(len(rows)))
    ax.set_yticklabels(rows, fontsize=11)
    ax.set_xlabel("cycle length  k", fontsize=10)
    for i in range(mat.shape[0]):
        for j in range(mat.shape[1]):
            v = mat[i, j]
            ax.text(j, i, f"{int(v)}", ha="center", va="center",
                    fontsize=12, fontweight="bold" if v else "normal",
                    color="white" if v > vmax * 0.5 else "#6b6b6b")
    # White cell separators (clean matrix look); hide the axes frame.
    ax.set_xticks(np.arange(-0.5, len(cats), 1), minor=True)
    ax.set_yticks(np.arange(-0.5, len(rows), 1), minor=True)
    ax.grid(which="minor", color="white", linewidth=2)
    ax.tick_params(which="minor", length=0)
    for sp in ax.spines.values():
        sp.set_visible(False)
    fig.tight_layout()
    fig.savefig(out_path, transparent=True, bbox_inches="tight")
    plt.close(fig)
    return True


def rasterize_chart_on(prs, marker: str, out_name: str) -> bool:  # type: ignore[no-untyped-def]
    """Replace the native (embedded) chart on the slide whose text contains
    `marker` with a pre-rendered PNG at the same box. Native charts depend on an
    embedded workbook and render blank in non-PowerPoint viewers / PDF export;
    a baked PNG renders identically everywhere. Returns True if one was swapped."""
    for slide in prs.slides:
        if marker not in slide_text(slide).lower():
            continue
        for sh in list(slide.shapes):
            if sh.has_chart:
                left, top, w, h = sh.left, sh.top, sh.width, sh.height
                png = FIG / out_name
                if not _render_cycle_spectrum(sh.chart, str(png)):
                    return False
                sh._element.getparent().remove(sh._element)
                slide.shapes.add_picture(str(png), left, top, width=w, height=h)
                return True
    return False


def consolidate_hymeyolo(prs) -> int:  # type: ignore[no-untyped-def]
    """Kato review: HyMeYOLO is "not today's main claim" — fold its three slides
    into ONE. The detections slide carries the figures and survives; the
    convolution-intro and Ricci/Hodge text slides (the "dragon heads" Kato wants
    tamed) are dropped. Returns the number of slides removed.

    Drop-only and two-phase (scan, then remove) — no in-place title mutation
    (which would reset fonts) and no removal during iteration (which would
    misindex the live slide view)."""
    sld_lst = prs.slides._sldIdLst
    ids = list(sld_lst)
    drop_markers = ("hypergraph convolution in vision", "ricci curvature")
    to_remove = [
        ids[i]
        for i, slide in enumerate(prs.slides)
        if any(m in slide_text(slide).lower() for m in drop_markers)
    ]
    for el in to_remove:
        # Drop the presentation→slide relationship too, not just the sldId entry,
        # so the package has no dangling rel (avoids the duplicate-part warning).
        prs.part.drop_rel(el.rId)
        sld_lst.remove(el)
    return len(to_remove)


def reorder_conclusion_last(prs) -> bool:  # type: ignore[no-untyped-def]
    """Move the 'Conclusion & outlook' slide to the very end (it must be last)."""
    sld_lst = prs.slides._sldIdLst
    ids = list(sld_lst)
    for i, slide in enumerate(prs.slides):
        # substring (not startswith): shape iteration order may place body text
        # before the title, so the joined text need not start with "Conclusion".
        if "conclusion" in slide_text(slide).lower():
            sld_lst.remove(ids[i])
            sld_lst.append(ids[i])
            return True
    return False


def add_bullet_slide(prs, title: str, bullets: list[str], body_pt: int) -> None:  # type: ignore[no-untyped-def]
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank)
    W, H = prs.slide_width, prs.slide_height
    tb = slide.shapes.add_textbox(Emu(int(W * 0.06)), Emu(int(H * 0.05)), Emu(int(W * 0.88)), Emu(int(H * 0.12)))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = title
    run.font.name = HEADER_FONT
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    body = slide.shapes.add_textbox(Emu(int(W * 0.06)), Emu(int(H * 0.20)), Emu(int(W * 0.88)), Emu(int(H * 0.74)))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = "•  " + line
        r.font.name = BODY_FONT
        r.font.size = Pt(body_pt)
        p.space_after = Pt(6)
    return slide


def move_to(prs, from_idx: int, to_idx: int) -> None:  # type: ignore[no-untyped-def]
    """Move a slide within the deck order (operates on the sldIdLst)."""
    sld_lst = prs.slides._sldIdLst
    ids = list(sld_lst)
    el = ids[from_idx]
    sld_lst.remove(el)
    sld_lst.insert(to_idx, el)


def add_publications_slide(prs) -> bool:  # type: ignore[no-untyped-def]
    """Publications & submissions: the portfolio figure + a grouped venue list."""
    fig = FIG / "publications.png"
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank)
    W, H = prs.slide_width, prs.slide_height
    tb = slide.shapes.add_textbox(Emu(int(W * 0.06)), Emu(int(H * 0.04)), Emu(int(W * 0.88)), Emu(int(H * 0.10)))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "Publications & submissions (2026)"
    r.font.name = HEADER_FONT
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    if fig.exists():
        slide.shapes.add_picture(str(fig), Emu(int(W * 0.30)), Emu(int(H * 0.15)), width=Emu(int(W * 0.66)))
    body = slide.shapes.add_textbox(Emu(int(W * 0.04)), Emu(int(H * 0.16)), Emu(int(W * 0.27)), Emu(int(H * 0.80)))
    tf = body.text_frame
    tf.word_wrap = True
    _pub_colour = {
        "sub": RGBColor(0x1b, 0x6c, 0xa8),
        "prep": RGBColor(0xc9, 0x7a, 0x16),
        "plan": RGBColor(0x70, 0x7a, 0x86),
    }
    for i, (line, tag) in enumerate(PUBLICATIONS):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = _PUB_PREFIX[tag] + line
        run.font.name = BODY_FONT
        run.font.size = Pt(9)
        run.font.color.rgb = _pub_colour[tag]
        p.space_after = Pt(4)
    return fig.exists()


def main() -> int:
    if not SRC.exists():
        print(f"missing deck: {SRC}")
        return 1
    prs = Presentation(str(SRC))
    placed = add_tensor_figure(prs)

    # Rasterize the fragile native column chart on "From structure to inductive
    # prior" (slide 18) to a clean PNG so it renders everywhere.
    chart_png = rasterize_chart_on(prs, "from structure to inductive prior", "kcycle_chart.png")

    # Compress (Kato + "merge slides"): HyMeYOLO 3 slides -> 1.
    dropped = consolidate_hymeyolo(prs)

    # Front: ONE consolidated contributions + arc + Gömb-gloss slide.
    add_bullet_slide(prs, "HyMeKo: contributions & talk arc", CONTRIB_ARC, body_pt=15)
    move_to(prs, -1, 1)

    add_bullet_slide(prs, "References (selected)", REFERENCES, body_pt=15)
    pubs = add_publications_slide(prs)
    add_bullet_slide(prs, "Future work", FUTURE_WORK, body_pt=17)

    # Back: ONE consolidated demo + collaboration slide, then anticipated Q&A.
    add_bullet_slide(prs, "Minimum demo & collaboration", DEMO_COLLAB, body_pt=16)
    add_bullet_slide(prs, "Anticipated questions (Kato)", QA, body_pt=13)

    concl = reorder_conclusion_last(prs)
    prs.save(str(DST))
    print(f"tensor background (square, faded) placed: {placed}")
    print(f"slide-18 k-cycle chart rasterized to PNG: {chart_png}")
    print(f"HyMeYOLO slides merged 3->1 (dropped {dropped})")
    print(f"publications figure embedded: {pubs}")
    print("added (front): contributions & talk arc (1 slide)")
    print("added: References + Publications + Future work")
    print("added (back): demo & collaboration (1) + anticipated Q&A (1)")
    print(f"conclusion moved to last: {concl}")
    print(f"total slides: {len(prs.slides._sldIdLst)}")
    print(f"wrote {DST}  (original {SRC.name} untouched)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
